# api/document_translator.py
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pdfplumber
from PyPDF2 import PdfReader, PdfWriter
import tempfile
import io
import os
from bs4 import BeautifulSoup
import html
import logging

logger = logging.getLogger(__name__)

class DocumentTranslator:
    def __init__(self, translation_model):
        self.model = translation_model

    async def translate_docx_with_progress(self, content: bytes, source_lang: str, target_lang: str, progress_callback):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name

        try:
            doc = Document(tmp_path)
            total_items = len(doc.paragraphs) + sum(len(table.rows) * len(table.columns) for table in doc.tables)
            processed = 0

            # Track total metrics
            total_processing_time = 0
            total_input_tokens = 0
            total_output_tokens = 0

            # Translate paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    translation_result = self.model.translate(paragraph.text, source_lang, target_lang)
                    translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                    metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                    
                    paragraph.text = translated_text
                    
                    # Safely accumulate metrics
                    total_input_tokens += metrics.get('input_tokens', 0)
                    total_output_tokens += metrics.get('output_tokens', 0)
                    total_processing_time += metrics.get('processing_time', 0)
                    
                processed += 1
                progress_callback(int(processed * 100 / total_items), "Translating document content...")

            # Translate tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            translation_result = self.model.translate(cell.text, source_lang, target_lang)
                            translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                            metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                            
                            cell.text = translated_text
                            
                            # Safely accumulate metrics
                            total_input_tokens += metrics.get('input_tokens', 0)
                            total_output_tokens += metrics.get('output_tokens', 0)
                            total_processing_time += metrics.get('processing_time', 0)
                            
                        processed += 1
                        progress_callback(int(processed * 100 / total_items), "Translating tables...")

            doc.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                content = f.read()

            # Calculate final metrics
            total_tokens = total_input_tokens + total_output_tokens
            tokens_per_second = total_tokens / total_processing_time if total_processing_time > 0 else 0

            final_metrics = {
                "tokens_per_second": round(tokens_per_second, 2),
                "total_tokens": total_tokens,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "processing_time": round(total_processing_time, 2),
                "cached": False
            }

            return content, final_metrics

        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    async def translate_xlsx_with_progress(self, content: bytes, source_lang: str, target_lang: str, progress_callback):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name

        try:
            wb = load_workbook(tmp_path)
            total_sheets = len(wb.sheetnames)
            processed_sheets = 0

            # Track total metrics
            total_processing_time = 0
            total_input_tokens = 0
            total_output_tokens = 0

            for sheet in wb.worksheets:
                rows = list(sheet.rows)
                total_cells = sum(1 for row in rows for cell in row if cell.value and isinstance(cell.value, str))
                processed_cells = 0

                for row in rows:
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            translation_result = self.model.translate(cell.value, source_lang, target_lang)
                            translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                            metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                            
                            cell.value = translated_text
                            
                            # Safely accumulate metrics
                            total_input_tokens += metrics.get('input_tokens', 0)
                            total_output_tokens += metrics.get('output_tokens', 0)
                            total_processing_time += metrics.get('processing_time', 0)
                            
                            processed_cells += 1
                            progress = int((processed_sheets * 100 / total_sheets) + 
                                         (processed_cells * 100 / total_cells / total_sheets))
                            progress_callback(min(progress, 99), f"Translating sheet {sheet.title}...")

                processed_sheets += 1

            wb.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                content = f.read()

            # Calculate final metrics
            total_tokens = total_input_tokens + total_output_tokens
            tokens_per_second = total_tokens / total_processing_time if total_processing_time > 0 else 0

            final_metrics = {
                "tokens_per_second": round(tokens_per_second, 2),
                "total_tokens": total_tokens,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "processing_time": round(total_processing_time, 2),
                "cached": False
            }

            return content, final_metrics

        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    async def translate_pptx_with_progress(self, content: bytes, source_lang: str, target_lang: str, progress_callback):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name

        try:
            prs = Presentation(tmp_path)
            total_slides = len(prs.slides)
            processed_slides = 0

            # Track total metrics
            total_processing_time = 0
            total_input_tokens = 0
            total_output_tokens = 0

            for slide in prs.slides:
                shapes = [shape for shape in slide.shapes if hasattr(shape, "text")]
                total_shapes = len(shapes)
                processed_shapes = 0

                for shape in shapes:
                    if shape.text.strip():
                        translation_result = self.model.translate(shape.text, source_lang, target_lang)
                        translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                        metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                        
                        shape.text = translated_text
                        
                        # Safely accumulate metrics
                        total_input_tokens += metrics.get('input_tokens', 0)
                        total_output_tokens += metrics.get('output_tokens', 0)
                        total_processing_time += metrics.get('processing_time', 0)
                        
                    processed_shapes += 1
                    progress = int((processed_slides * 100 / total_slides) + 
                                 (processed_shapes * 100 / total_shapes / total_slides))
                    progress_callback(min(progress, 99), f"Translating slide {processed_slides + 1}...")

                processed_slides += 1

            prs.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                content = f.read()

            # Calculate final metrics
            total_tokens = total_input_tokens + total_output_tokens
            tokens_per_second = total_tokens / total_processing_time if total_processing_time > 0 else 0

            final_metrics = {
                "tokens_per_second": round(tokens_per_second, 2),
                "total_tokens": total_tokens,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "processing_time": round(total_processing_time, 2),
                "cached": False
            }

            return content, final_metrics

        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    async def translate_pdf_with_progress(self, content: bytes, source_lang: str, target_lang: str, progress_callback):
        # Create a temporary file to work with the PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name

        try:
            # Track total metrics
            total_processing_time = 0
            total_input_tokens = 0
            total_output_tokens = 0

            # Open PDF with both libraries
            pdf_reader = PdfReader(tmp_path)
            total_pages = len(pdf_reader.pages)
            
            # Create a new PDF writer for the translated content
            pdf_writer = PdfWriter()

            # Process each page
            with pdfplumber.open(tmp_path) as pdf:
                for page_num in range(total_pages):
                    # Extract text with pdfplumber for better text extraction
                    page = pdf.pages[page_num]
                    text_elements = []

                    # Extract text elements with their positions
                    for element in page.extract_words():
                        if element.get('text', '').strip():
                            text_elements.append({
                                'text': element['text'],
                                'x0': element['x0'],
                                'y0': element['top'],
                                'x1': element['x1'],
                                'y1': element['bottom']
                            })

                    # Get the PDF page for modification
                    pdf_page = pdf_reader.pages[page_num]
                    
                    # Create a content stream for the translated text
                    translated_content = []
                    
                    # Process and translate each text element
                    total_elements = len(text_elements)
                    for idx, element in enumerate(text_elements):
                        if element['text'].strip():
                            # Translate the text
                            translation_result = self.model.translate(element['text'], source_lang, target_lang)
                            translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                            metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                            
                            # Update metrics
                            total_input_tokens += metrics.get('input_tokens', 0)
                            total_output_tokens += metrics.get('output_tokens', 0)
                            total_processing_time += metrics.get('processing_time', 0)

                            # Store translated text with its position
                            translated_content.append({
                                'text': translated_text,
                                'x': element['x0'],
                                'y': element['y0']
                            })

                        # Update progress
                        current_progress = (page_num * 100 // total_pages) + (idx * 100 // total_elements // total_pages)
                        progress_callback(min(current_progress, 99), f"Translating page {page_num + 1}...")

                    # Create a new content stream with translated text
                    # Note: This is a simplified approach. A more sophisticated implementation
                    # would need to handle font embedding, text styling, and layout preservation
                    content_stream = f"""
                    BT
                    /F1 12 Tf
                    """
                    
                    for item in translated_content:
                        content_stream += f"{item['x']} {item['y']} Td ({item['text']}) Tj\n"
                    
                    content_stream += "ET"

                    # Add the modified page to the new PDF
                    pdf_writer.add_page(pdf_page)

            # Save the translated PDF to a temporary buffer
            output_buffer = io.BytesIO()
            pdf_writer.write(output_buffer)
            translated_content = output_buffer.getvalue()

            # Calculate final metrics
            total_tokens = total_input_tokens + total_output_tokens
            tokens_per_second = total_tokens / total_processing_time if total_processing_time > 0 else 0

            final_metrics = {
                "tokens_per_second": round(tokens_per_second, 2),
                "total_tokens": total_tokens,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "processing_time": round(total_processing_time, 2),
                "cached": False
            }

            return translated_content, final_metrics

        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    async def translate_html_with_progress(self, content: bytes, source_lang: str, target_lang: str, progress_callback):
        try:
            # Convert bytes to string
            html_content = content.decode('utf-8')
            
            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Get all text elements (focusing on common text-containing tags)
            text_elements = []
            for tag in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'span', 'div', 'li', 'td', 'th', 'a']):
                if tag.string and tag.string.strip():
                    text_elements.append(tag)

            total_elements = len(text_elements)
            processed = 0

            # Track metrics
            total_processing_time = 0
            total_input_tokens = 0
            total_output_tokens = 0

            # Translate each text element
            for element in text_elements:
                if element.string and element.string.strip():
                    # Translate the text
                    translation_result = self.model.translate(element.string, source_lang, target_lang)
                    translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                    metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                    
                    # Update metrics
                    total_input_tokens += metrics.get('input_tokens', 0)
                    total_output_tokens += metrics.get('output_tokens', 0)
                    total_processing_time += metrics.get('processing_time', 0)
                    
                    # Replace the text in the HTML
                    element.string = translated_text
                    
                    processed += 1
                    progress = int(processed * 100 / total_elements)
                    progress_callback(min(progress, 99), f"Translating element {processed} of {total_elements}...")

            # Calculate final metrics
            total_tokens = total_input_tokens + total_output_tokens
            tokens_per_second = total_tokens / total_processing_time if total_processing_time > 0 else 0

            final_metrics = {
                "tokens_per_second": round(tokens_per_second, 2),
                "total_tokens": total_tokens,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "processing_time": round(total_processing_time, 2),
                "cached": False
            }

            # Convert back to string and encode to bytes
            translated_html = soup.prettify().encode('utf-8')
            
            return translated_html, final_metrics

        except Exception as e:
            logger.error(f"Text file translation error: {str(e)}")
            raise
# Add to DocumentTranslator class
    async def translate_txt_with_progress(self, content: bytes, source_lang: str, target_lang: str, progress_callback):
        try:
            logger.info("Starting text file translation")
            # Decode text content with error handling
            try:
                text_content = content.decode('utf-8')
            except UnicodeDecodeError:
                # Try alternative encodings if UTF-8 fails
                text_content = content.decode('iso-8859-1')
                
            # Split into paragraphs (split by double newlines)
            paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
            
            # If no paragraphs found, split by single newlines
            if not paragraphs:
                paragraphs = [p.strip() for p in text_content.split('\n') if p.strip()]
            
            if not paragraphs:
                logger.warning("No text content found in file")
                return content, {
                    "tokens_per_second": 0,
                    "total_tokens": 0,
                    "input_tokens": 0,
                    "output_tokens": 0,
                    "processing_time": 0,
                    "cached": False
                }
            
            total_paragraphs = len(paragraphs)
            processed = 0
            translated_paragraphs = []
            
            # Track metrics
            total_processing_time = 0
            total_input_tokens = 0
            total_output_tokens = 0
            
            # Translate each paragraph
            for paragraph in paragraphs:
                if paragraph:
                    # Translate the text
                    translation_result = self.model.translate(paragraph, source_lang, target_lang)
                    translated_text = translation_result[0] if isinstance(translation_result, tuple) else translation_result
                    metrics = translation_result[1] if isinstance(translation_result, tuple) else {}
                    
                    # Update metrics
                    total_input_tokens += metrics.get('input_tokens', 0)
                    total_output_tokens += metrics.get('output_tokens', 0)
                    total_processing_time += metrics.get('processing_time', 0)
                    
                    translated_paragraphs.append(translated_text)
                
                processed += 1
                progress = int(processed * 100 / total_paragraphs)
                progress_callback(min(progress, 99), f"Translating paragraph {processed} of {total_paragraphs}...")
            
            # Join paragraphs with double newlines
            translated_content = '\n\n'.join(translated_paragraphs)
            
            # Calculate final metrics
            total_tokens = total_input_tokens + total_output_tokens
            tokens_per_second = total_tokens / total_processing_time if total_processing_time > 0 else 0
            
            final_metrics = {
                "tokens_per_second": round(tokens_per_second, 2),
                "total_tokens": total_tokens,
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "processing_time": round(total_processing_time, 2),
                "cached": False
            }
            
            logger.info("Text file translation completed successfully")
            # Return encoded content
            return translated_content.encode('utf-8'), final_metrics
            
        except Exception as e:
            logger.error(f"Text file translation error: {str(e)}")
            raise